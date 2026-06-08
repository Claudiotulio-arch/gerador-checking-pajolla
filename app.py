import io
import os
import re
import tempfile
from copy import deepcopy
from pathlib import Path

import streamlit as st
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Pt

APP_DIR = Path(__file__).parent
TEMPLATE_PATH = APP_DIR / "template_checking_pajolla.pptx"

st.set_page_config(page_title="Gerador de Checking Pajolla", page_icon="📸", layout="wide")

st.markdown(
    """
    <style>
      .main {background: #f7f8fb;}
      .block-container {padding-top: 2rem; max-width: 1180px;}
      h1, h2, h3 {color: #1f2a44;}
      .stButton button, .stDownloadButton button {
        border-radius: 10px;
        font-weight: 700;
      }
      .info-card {
        padding: 18px 20px;
        background: white;
        border-radius: 14px;
        border: 1px solid #e8eaf1;
        box-shadow: 0 1px 8px rgba(20, 30, 60, .05);
      }
    </style>
    """,
    unsafe_allow_html=True,
)

st.title("Gerador de Checking Pajolla")
st.caption("Preencha os campos, envie as fotos e baixe o PowerPoint pronto. Versão corrigida 1.1")


def _delete_slide(prs: Presentation, index: int) -> None:
    """Remove a slide by index from a Presentation."""
    xml_slides = prs.slides._sldIdLst  # noqa: SLF001
    slides = list(xml_slides)
    r_id = slides[index].rId
    prs.part.drop_rel(r_id)
    xml_slides.remove(slides[index])


def _duplicate_slide(prs: Presentation, source_index: int):
    """Duplicate a slide by copying shapes and relationships into a blank slide.

    This avoids an incompatibility in newer python-pptx versions where directly
    moving slide XML can raise: NoneType object has no attribute remove.
    """
    source = prs.slides[source_index]
    new_slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Remove any placeholders created by the blank layout.
    for shp in list(new_slide.shapes):
        new_slide.shapes._spTree.remove(shp._element)  # noqa: SLF001

    # Copy relationships used by shapes, mapping old rIds to the new slide rIds.
    rel_id_map = {}
    for rel in source.part.rels.values():
        if "notesSlide" in rel.reltype or "slideLayout" in rel.reltype:
            continue
        if rel.is_external:
            new_rid = new_slide.part.relate_to(rel.target_ref, rel.reltype, is_external=True)
        else:
            new_rid = new_slide.part.relate_to(rel.target_part, rel.reltype)
        rel_id_map[rel.rId] = new_rid

    # Copy all shapes and update any relationship ids inside the copied XML.
    for shp in source.shapes:
        new_el = deepcopy(shp._element)  # noqa: SLF001
        for el in new_el.iter():
            for attr_name, attr_value in list(el.attrib.items()):
                if attr_value in rel_id_map:
                    el.attrib[attr_name] = rel_id_map[attr_value]
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")  # noqa: SLF001

    return new_slide


def _all_text(shape):
    if hasattr(shape, "text"):
        return shape.text or ""
    return ""


def _replace_text_preserving_style(shape, new_text: str) -> None:
    """Replace text while preserving most paragraph/run formatting."""
    if not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    if not tf.paragraphs:
        shape.text = new_text
        return
    first_para = tf.paragraphs[0]
    if not first_para.runs:
        first_para.text = new_text
        return
    first_run = first_para.runs[0]
    first_run.text = new_text
    for para in tf.paragraphs:
        for run in para.runs[1:]:
            run.text = ""
    for para in tf.paragraphs[1:]:
        para.text = ""


def _replace_exact_text(slide, replacements: dict) -> None:
    for shape in slide.shapes:
        text = _all_text(shape).strip()
        if text in replacements:
            _replace_text_preserving_style(shape, replacements[text])


def _replace_by_position(slide, values: dict) -> None:
    """Update checking fields by using the template labels and nearby value shapes."""
    labels = {}
    for shape in slide.shapes:
        text = _all_text(shape).strip().lower()
        if text in {"cliente", "mídia", "midia", "local", "atendimento", "período", "periodo"}:
            labels[text] = shape

    def set_right_of(label_keys, value):
        label_shape = None
        for key in label_keys:
            if key in labels:
                label_shape = labels[key]
                break
        if not label_shape:
            return
        candidates = []
        for shp in slide.shapes:
            if shp == label_shape or not getattr(shp, "has_text_frame", False):
                continue
            if shp.left <= label_shape.left:
                continue
            vertical_distance = abs((shp.top + shp.height / 2) - (label_shape.top + label_shape.height / 2))
            if vertical_distance < Pt(28):
                candidates.append((vertical_distance, shp.left, shp))
        if candidates:
            _, _, shp = sorted(candidates, key=lambda x: (x[0], x[1]))[0]
            _replace_text_preserving_style(shp, value)

    set_right_of(["cliente"], values["cliente"])
    set_right_of(["mídia", "midia"], values["midia"])
    set_right_of(["local"], values["local"])
    set_right_of(["atendimento"], values["atendimento"])
    set_right_of(["período", "periodo"], values["periodo_campanha"])

    # Top period box. It is usually the shape starting with PERÍODO: or the month text.
    for shp in slide.shapes:
        text = _all_text(shp).strip()
        if text.startswith("PERÍODO:") or re.search(r"\b(ABR|MAI|JUN|JUL|AGO|SET|OUT|NOV|DEZ|JAN|FEV|MAR)", text.upper()):
            if shp.left < 3_500_000 and shp.top < 1_300_000:
                _replace_text_preserving_style(shp, f"PERÍODO: {values['periodo_topo']}")


def _find_photo_placeholder(slide):
    # Prefer the shape with the text FOTO. If not found, use the largest shape on the right side.
    for shp in slide.shapes:
        if _all_text(shp).strip().upper() == "FOTO":
            return shp
    candidates = [shp for shp in slide.shapes if shp.left > 2_500_000]
    if not candidates:
        return None
    return max(candidates, key=lambda s: s.width * s.height)


def _insert_image_no_distortion(slide, image_path: str, box, fill_area: bool) -> None:
    """Insert image in the placeholder area without distortion. Can contain or cover."""
    with Image.open(image_path) as im:
        img_w, img_h = im.size
    img_ratio = img_w / img_h
    box_ratio = box.width / box.height

    if fill_area:
        # Cover: fills the photo area, but may crop due to PowerPoint's internal scaling.
        slide.shapes.add_picture(image_path, box.left, box.top, width=box.width, height=box.height)
        return

    # Contain: preserves full image, centered, no crop.
    if img_ratio > box_ratio:
        new_w = box.width
        new_h = int(box.width / img_ratio)
        left = box.left
        top = int(box.top + (box.height - new_h) / 2)
    else:
        new_h = box.height
        new_w = int(box.height * img_ratio)
        left = int(box.left + (box.width - new_w) / 2)
        top = box.top

    # Put a white background behind the image so the word FOTO doesn't appear around it.
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, box.left, box.top, box.width, box.height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.fill.background()
    slide.shapes.add_picture(image_path, left, top, width=new_w, height=new_h)


def _safe_filename(text: str) -> str:
    text = re.sub(r"[^A-Za-z0-9À-ÿ_-]+", "_", text.strip())
    text = re.sub(r"_+", "_", text).strip("_")
    return text or "checking"


def generate_pptx(template_path: Path, photo_files, fields: dict, fill_area: bool) -> bytes:
    prs = Presentation(str(template_path))

    # Keep the cover and the checking slides already present in the model.
    # Extra slides are deleted after we know how many photos were uploaded.

    # Update cover slide fields, preserving the model's original look.
    if len(prs.slides) >= 1:
        _replace_exact_text(
            prs.slides[0],
            {
                "BYD": fields["cliente"],
                "2984": fields.get("contrato", ""),
                "MÍDIA OFF": fields.get("campanha", "MÍDIA OFF"),
            },
        )

    temp_paths = []
    try:
        # Save uploads into temp files because python-pptx needs paths.
        for idx, file in enumerate(photo_files, start=1):
            suffix = Path(file.name).suffix or ".jpg"
            tmp = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
            tmp.write(file.getvalue())
            tmp.close()
            temp_paths.append(tmp.name)

        # Use the checking slides already in the model. Duplicate the first checking
        # slide only if the user uploads more photos than the template contains.
        while len(prs.slides) < len(temp_paths) + 1:
            _duplicate_slide(prs, 1)
        while len(prs.slides) > len(temp_paths) + 1:
            _delete_slide(prs, len(prs.slides) - 1)

        for i, image_path in enumerate(temp_paths, start=1):
            slide = prs.slides[i]
            _replace_by_position(slide, fields)
            placeholder = _find_photo_placeholder(slide)
            if placeholder is not None:
                _insert_image_no_distortion(slide, image_path, placeholder, fill_area=fill_area)

        out = io.BytesIO()
        prs.save(out)
        return out.getvalue()
    finally:
        for p in temp_paths:
            try:
                os.unlink(p)
            except OSError:
                pass


with st.sidebar:
    st.header("Campos da campanha")
    cliente = st.text_input("Cliente", value="COLÉGIO POSITIVO")
    periodo_topo = st.text_input("Período no topo", value="JUN/ 2026")
    midia = st.text_input("Mídia", value="SURF CENTER")
    local = st.text_input("Local", value="CURITIBA")
    atendimento = st.text_input("Atendimento", value="SHINOHARA")
    periodo_campanha = st.text_input("Período da campanha", value="JUNHO/2026")
    contrato = st.text_input("Nº do Contrato/P.I", value="")
    campanha = st.text_input("Campanha", value="MÍDIA OFF")
    modo = st.radio(
        "Como inserir a foto?",
        ["Foto inteira, sem corte", "Preencher todo o espaço"],
        index=0,
        help="Foto inteira não corta nem deforma. Preencher todo o espaço pode cortar bordas, mas ocupa 100% da área.",
    )

st.markdown('<div class="info-card">', unsafe_allow_html=True)
st.subheader("1. Envie as fotos")
photos = st.file_uploader(
    "Selecione uma ou várias fotos do checking",
    type=["jpg", "jpeg", "png", "webp"],
    accept_multiple_files=True,
)
st.markdown("</div>", unsafe_allow_html=True)

if photos:
    st.subheader("Prévia das fotos enviadas")
    cols = st.columns(4)
    for idx, photo in enumerate(photos):
        with cols[idx % 4]:
            st.image(photo, caption=photo.name, use_container_width=True)

fields = {
    "cliente": cliente.upper().strip(),
    "periodo_topo": periodo_topo.upper().strip(),
    "midia": midia.upper().strip(),
    "local": local.upper().strip(),
    "atendimento": atendimento.upper().strip(),
    "periodo_campanha": periodo_campanha.strip(),
    "contrato": contrato.strip(),
    "campanha": campanha.upper().strip(),
}

st.divider()

col1, col2 = st.columns([1, 2])
with col1:
    gerar = st.button("Gerar PPTX", type="primary", disabled=not photos)
with col2:
    if not photos:
        st.info("Envie pelo menos uma foto para gerar o checking.")

if gerar:
    if not TEMPLATE_PATH.exists():
        st.error("Arquivo template_checking_pajolla.pptx não encontrado na pasta do app.")
    else:
        with st.spinner("Gerando PowerPoint..."):
            pptx_bytes = generate_pptx(
                TEMPLATE_PATH,
                photos,
                fields,
                fill_area=(modo == "Preencher todo o espaço"),
            )
        file_name = f"checking_{_safe_filename(cliente)}_{_safe_filename(periodo_topo)}.pptx"
        st.success("Checking gerado com sucesso.")
        st.download_button(
            "Baixar PPTX",
            data=pptx_bytes,
            file_name=file_name,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            type="primary",
        )
